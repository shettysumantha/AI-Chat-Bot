import os
import re
from typing import List


class DocumentParser:
    @staticmethod
    def extract_text(file_path: str, file_type: str) -> str:
        ext = (file_type or os.path.splitext(file_path)[1] or '').lower()
        if ext == '.txt' or ext == '.md':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as handle:
                return handle.read()
        if ext == '.csv':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as handle:
                return handle.read()
        if ext == '.pdf':
            try:
                import PyPDF2
            except Exception:
                return ''
            reader = PyPDF2.PdfReader(file_path)
            text = '\n'.join(page.extract_text() or '' for page in reader.pages)
            return text
        if ext == '.docx':
            try:
                import docx
            except Exception:
                return ''
            document = docx.Document(file_path)
            return '\n'.join(p.text for p in document.paragraphs)
        if ext == '.xlsx':
            try:
                import openpyxl
            except Exception:
                return ''
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            sheets = []
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    sheets.append(' | '.join(str(cell) for cell in row if cell is not None))
            return '\n'.join(sheets)
        if ext == '.pptx':
            try:
                from pptx import Presentation
            except Exception:
                return ''
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        texts.append(shape.text)
            return '\n'.join(texts)
        return ''

    @staticmethod
    def split_text(text: str, chunk_size: int = 600, overlap: int = 120) -> List[str]:
        if not text:
            return []
        cleaned = re.sub(r'\s+', ' ', text).strip()
        if not cleaned:
            return []
        chunks = []
        start = 0
        while start < len(cleaned):
            end = start + chunk_size
            chunk = cleaned[start:end]
            if not chunk:
                break
            chunks.append(chunk)
            if end >= len(cleaned):
                break
            start += chunk_size - overlap
        return chunks
